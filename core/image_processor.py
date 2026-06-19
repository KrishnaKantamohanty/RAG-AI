import io
import os
import base64
from PIL import Image
import pytesseract
import fitz  # PyMuPDF
import docx
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE

class ImageProcessor:
    """
    ImageProcessor handles image processing, text extraction via OCR,
    and image extraction from documents (PDF, DOCX, PPTX).
    """
    def __init__(self):
        tess_path = os.getenv("TESSERACT_PATH")
        if tess_path and os.path.exists(tess_path):
            pytesseract.pytesseract.tesseract_cmd = tess_path
            
        self.tesseract_available = False
        try:
            pytesseract.get_tesseract_version()
            self.tesseract_available = True
        except Exception:
            print("OCR Warning: Tesseract is not installed or not in PATH. OCR features will be disabled.")

    def run_ocr(self, pil_image: Image.Image) -> str:
        if not self.tesseract_available:
            return ""
        try:
            text = pytesseract.image_to_string(pil_image, lang='eng')
            # Clean up empty lines
            lines = [line.strip() for line in text.split('\n') if line.strip()]
            return "\n".join(lines)
        except Exception as e:
            print(f"OCR Error: {e}")
            return ""

    def image_to_base64(self, pil_image: Image.Image) -> str:
        img = pil_image.copy()
        if img.mode != "RGB":
            img = img.convert("RGB")
            
        # Resize if larger than 2000px on any side
        max_size = 2000
        if img.width > max_size or img.height > max_size:
            img.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)
            
        buffer = io.BytesIO()
        # Compress to JPEG
        img.save(buffer, format="JPEG", quality=85)
        
        # Check size if > 2MB and compress more if needed
        if buffer.tell() > 2 * 1024 * 1024:
            buffer = io.BytesIO()
            img.save(buffer, format="JPEG", quality=60)
            
        return base64.b64encode(buffer.getvalue()).decode("utf-8")

    def process_image_file(self, uploaded_file) -> dict:
        filename = getattr(uploaded_file, "filename", getattr(uploaded_file, "name", "unknown_image"))
        file_bytes = uploaded_file.read()
        uploaded_file.seek(0)  # Reset pointer just in case
        
        try:
            img = Image.open(io.BytesIO(file_bytes))
            if img.mode != "RGB":
                img = img.convert("RGB")
                
            b64_str = self.image_to_base64(img)
            ocr_text = self.run_ocr(img)
            
            return {
                "image": img,
                "base64": b64_str,
                "filename": filename,
                "size": img.size,
                "mode": img.mode,
                "ocr_text": ocr_text,
                "type": "image"
            }
        except Exception as e:
            print(f"Error processing image file: {e}")
            return {}

    def extract_images_from_pdf(self, uploaded_file) -> list[dict]:
        images_data = []
        file_bytes = uploaded_file.read()
        uploaded_file.seek(0)
        
        try:
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # Render full page
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                if img.mode != "RGB":
                    img = img.convert("RGB")
                    
                b64_str = self.image_to_base64(img)
                ocr_text = self.run_ocr(img)
                
                images_data.append({
                    "image": img,
                    "base64": b64_str,
                    "page_number": page_num + 1,
                    "ocr_text": ocr_text,
                    "type": "pdf_page"
                })
                
                # Also extract embedded images
                image_list = page.get_images(full=True)
                for img_index, img_info in enumerate(image_list):
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    try:
                        emb_img = Image.open(io.BytesIO(image_bytes))
                        if emb_img.mode != "RGB":
                            emb_img = emb_img.convert("RGB")
                        
                        # Only keep somewhat large images to avoid tiny icons
                        if emb_img.width > 100 and emb_img.height > 100:
                            emb_b64 = self.image_to_base64(emb_img)
                            emb_ocr = self.run_ocr(emb_img)
                            
                            images_data.append({
                                "image": emb_img,
                                "base64": emb_b64,
                                "page_number": page_num + 1,
                                "ocr_text": emb_ocr,
                                "type": "pdf_embedded_image"
                            })
                    except Exception as e:
                        print(f"Failed to process embedded pdf image: {e}")
                        
            return images_data
        except Exception as e:
            print(f"Error extracting images from PDF: {e}")
            return []

    def extract_images_from_docx(self, uploaded_file) -> list[dict]:
        images_data = []
        file_bytes = uploaded_file.read()
        uploaded_file.seek(0)
        
        try:
            doc_file = docx.Document(io.BytesIO(file_bytes))
            for rel in doc_file.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_bytes = rel.target_part.blob
                        img = Image.open(io.BytesIO(image_bytes))
                        if img.mode != "RGB":
                            img = img.convert("RGB")
                            
                        if img.width > 100 and img.height > 100:
                            b64_str = self.image_to_base64(img)
                            ocr_text = self.run_ocr(img)
                            
                            images_data.append({
                                "image": img,
                                "base64": b64_str,
                                "ocr_text": ocr_text,
                                "type": "docx_image"
                            })
                    except Exception as e:
                        print(f"Failed to process docx image: {e}")
            return images_data
        except Exception as e:
            print(f"Error extracting images from DOCX: {e}")
            return []

    def extract_images_from_pptx(self, uploaded_file) -> list[dict]:
        images_data = []
        file_bytes = uploaded_file.read()
        uploaded_file.seek(0)
        
        try:
            prs = pptx.Presentation(io.BytesIO(file_bytes))
            for slide_num, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image_bytes = shape.image.blob
                            img = Image.open(io.BytesIO(image_bytes))
                            if img.mode != "RGB":
                                img = img.convert("RGB")
                                
                            if img.width > 100 and img.height > 100:
                                b64_str = self.image_to_base64(img)
                                ocr_text = self.run_ocr(img)
                                
                                images_data.append({
                                    "slide_number": slide_num + 1,
                                    "base64": b64_str,
                                    "ocr_text": ocr_text,
                                    "type": "pptx_image"
                                })
                        except Exception as e:
                            print(f"Failed to process pptx image: {e}")
            return images_data
        except Exception as e:
            print(f"Error extracting images from PPTX: {e}")
            return []
