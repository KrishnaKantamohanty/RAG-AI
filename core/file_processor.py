import pandas as pd
import pdfplumber
import docx
import pptx
import json
import csv
import io
import os

from core.image_processor import ImageProcessor

class FileProcessor:
    """
    FileProcessor handles extraction of text and images from various file formats.
    It supports PDF, Word, Excel, PowerPoint, TXT, CSV, JSON, and raw image files.
    """
    
    def process(self, file) -> tuple[str, list[dict]]:
        """
        Accepts a Werkzeug FileStorage and returns extracted text and a list of image dicts.
        Returns: (text: str, images: list[dict])
        """
        try:
            filename = getattr(file, "filename", getattr(file, "name", "unknown")).lower()
            file_bytes = file.read()
            file.seek(0)
            
            extracted_text = []
            extracted_images = []
            image_processor = ImageProcessor()
            
            image_extensions = (".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp", ".tiff")
            if filename.endswith(image_extensions):
                img_data = image_processor.process_image_file(file)
                if img_data:
                    extracted_images.append(img_data)
                return ("", extracted_images)
            
            if filename.endswith(".pdf"):
                with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            extracted_text.append(text)
                
                # Rewind and extract images
                file.seek(0)
                extracted_images.extend(image_processor.extract_images_from_pdf(file))
                            
            elif filename.endswith(".docx") or filename.endswith(".doc"):
                doc_obj = docx.Document(io.BytesIO(file_bytes))
                for para in doc_obj.paragraphs:
                    if para.text.strip():
                        extracted_text.append(para.text)
                for table in doc_obj.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                extracted_text.append(cell.text)
                                
                file.seek(0)
                extracted_images.extend(image_processor.extract_images_from_docx(file))
                                
            elif filename.endswith((".xlsx", ".xls")):
                excel_file = pd.ExcelFile(io.BytesIO(file_bytes))
                for sheet in excel_file.sheet_names:
                    df = excel_file.parse(sheet)
                    extracted_text.append(f"Sheet: {sheet}")
                    extracted_text.append(df.to_markdown(index=False))
                    
            elif filename.endswith(".csv"):
                df = pd.read_csv(io.BytesIO(file_bytes))
                extracted_text.append(df.to_markdown(index=False))
                
            elif filename.endswith((".pptx", ".ppt")):
                prs = pptx.Presentation(io.BytesIO(file_bytes))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            extracted_text.append(shape.text)
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                        if slide.notes_slide.notes_text_frame.text.strip():
                            extracted_text.append(slide.notes_slide.notes_text_frame.text)
                            
                file.seek(0)
                extracted_images.extend(image_processor.extract_images_from_pptx(file))
                            
            elif filename.endswith(".txt"):
                extracted_text.append(file_bytes.decode('utf-8'))
                
            elif filename.endswith(".json"):
                data = json.loads(file_bytes.decode('utf-8'))
                extracted_text.append(json.dumps(data, indent=2))
                
            else:
                return (f"Error: Unsupported file type for {filename}", [])
                
            return ("\n\n".join(extracted_text), extracted_images)
            
        except Exception as e:
            return (f"Error extracting content from {getattr(file, 'filename', 'file')}: {str(e)}", [])
